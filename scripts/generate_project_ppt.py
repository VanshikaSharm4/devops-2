#!/usr/bin/env python3
"""Generate DevOps Agent 2 project overview PowerPoint."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

# Adobe-inspired palette
BRAND_BLUE = RGBColor(0x14, 0x73, 0xE6)
DARK = RGBColor(0x1A, 0x1A, 0x1A)
GRAY = RGBColor(0x4B, 0x4B, 0x4B)
LIGHT_GRAY = RGBColor(0x90, 0x90, 0x90)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
RED = RGBColor(0xE5, 0x48, 0x4D)
GREEN = RGBColor(0x2D, 0x9D, 0x5C)
AMBER = RGBColor(0xC9, 0x89, 0x00)


def _set_bg(slide, color=WHITE):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_accent_bar(slide, top=Inches(0), height=Inches(0.08)):
    shape = slide.shapes.add_shape(
        1, Inches(0), top, Inches(13.33), height  # MSO_SHAPE.RECTANGLE = 1
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = BRAND_BLUE
    shape.line.fill.background()


def _title_slide(prs, title, subtitle, footer=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    _set_bg(slide)
    _add_accent_bar(slide, top=Inches(0), height=Inches(0.12))

    box = slide.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(11.5), Inches(1.5))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = DARK

    sub = slide.shapes.add_textbox(Inches(0.8), Inches(3.6), Inches(11.5), Inches(1.2))
    sp = sub.text_frame.paragraphs[0]
    sp.text = subtitle
    sp.font.size = Pt(20)
    sp.font.color.rgb = GRAY

    if footer:
        ft = slide.shapes.add_textbox(Inches(0.8), Inches(6.8), Inches(11.5), Inches(0.5))
        fp = ft.text_frame.paragraphs[0]
        fp.text = footer
        fp.font.size = Pt(14)
        fp.font.color.rgb = LIGHT_GRAY


def _section_slide(prs, section_num, title):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, RGBColor(0x14, 0x73, 0xE6))
    box = slide.shapes.add_textbox(Inches(0.8), Inches(2.8), Inches(11.5), Inches(1.5))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = f"{section_num:02d}"
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(0xBB, 0xD4, 0xF7)
    p2 = tf.add_paragraph()
    p2.text = title
    p2.font.size = Pt(36)
    p2.font.bold = True
    p2.font.color.rgb = WHITE


def _content_slide(prs, title, bullets, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide)
    _add_accent_bar(slide)

    tb = slide.shapes.add_textbox(Inches(0.7), Inches(0.35), Inches(12), Inches(0.7))
    tp = tb.text_frame.paragraphs[0]
    tp.text = title
    tp.font.size = Pt(28)
    tp.font.bold = True
    tp.font.color.rgb = DARK

    y = 1.15
    if subtitle:
        sb = slide.shapes.add_textbox(Inches(0.7), Inches(1.0), Inches(12), Inches(0.4))
        sp = sb.text_frame.paragraphs[0]
        sp.text = subtitle
        sp.font.size = Pt(14)
        sp.font.italic = True
        sp.font.color.rgb = LIGHT_GRAY
        y = 1.45

    body = slide.shapes.add_textbox(Inches(0.7), Inches(y), Inches(12), Inches(5.8))
    tf = body.text_frame
    tf.word_wrap = True
    for i, item in enumerate(bullets):
        if isinstance(item, tuple):
            text, level = item
        else:
            text, level = item, 0
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.level = level
        p.font.size = Pt(16 if level == 0 else 14)
        p.font.color.rgb = GRAY if level == 0 else LIGHT_GRAY
        p.space_after = Pt(8)


def _two_column_slide(prs, title, left_title, left_items, right_title, right_items):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide)
    _add_accent_bar(slide)

    tb = slide.shapes.add_textbox(Inches(0.7), Inches(0.35), Inches(12), Inches(0.7))
    tp = tb.text_frame.paragraphs[0]
    tp.text = title
    tp.font.size = Pt(28)
    tp.font.bold = True
    tp.font.color.rgb = DARK

    for col_x, col_title, items in [
        (0.7, left_title, left_items),
        (6.8, right_title, right_items),
    ]:
        ht = slide.shapes.add_textbox(Inches(col_x), Inches(1.2), Inches(5.8), Inches(0.4))
        hp = ht.text_frame.paragraphs[0]
        hp.text = col_title
        hp.font.size = Pt(16)
        hp.font.bold = True
        hp.font.color.rgb = BRAND_BLUE

        body = slide.shapes.add_textbox(Inches(col_x), Inches(1.65), Inches(5.8), Inches(5.2))
        tf = body.text_frame
        tf.word_wrap = True
        for i, text in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = f"• {text}"
            p.font.size = Pt(14)
            p.font.color.rgb = GRAY
            p.space_after = Pt(6)


def _table_slide(prs, title, headers, rows):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide)
    _add_accent_bar(slide)

    tb = slide.shapes.add_textbox(Inches(0.7), Inches(0.35), Inches(12), Inches(0.7))
    tp = tb.text_frame.paragraphs[0]
    tp.text = title
    tp.font.size = Pt(28)
    tp.font.bold = True
    tp.font.color.rgb = DARK

    n_rows = len(rows) + 1
    n_cols = len(headers)
    table_shape = slide.shapes.add_table(n_rows, n_cols, Inches(0.7), Inches(1.2), Inches(12), Inches(0.4 * n_rows))
    table = table_shape.table

    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.bold = True
            p.font.size = Pt(12)
            p.font.color.rgb = WHITE
        cell.fill.solid()
        cell.fill.fore_color.rgb = BRAND_BLUE

    for i, row in enumerate(rows, start=1):
        for j, val in enumerate(row):
            cell = table.cell(i, j)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(11)
                p.font.color.rgb = GRAY


def build_presentation() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ── Title ──
    _title_slide(
        prs,
        "DevOps Intelligence Platform",
        "AI-Powered Pipeline Intelligence for Adobe Cloud Manager",
        "IDFC First Bank · Program 19905 · Adobe Internal",
    )

    # ── Agenda ──
    _content_slide(
        prs,
        "Agenda",
        [
            "Problem Statement & Business Context",
            "Solution Overview & Value Proposition",
            "8-Layer Data Processing Architecture",
            "Core Features (6 Intelligence Capabilities)",
            "LogSage: Post-Failure Assessment Pipeline",
            "Hybrid Retrieval & Institutional Memory (RAG)",
            "ML Risk Integration",
            "Streamlit Dashboard & CLI",
            "Technology Stack & Current Status",
        ],
    )

    # ── Section 1: Problem ──
    _section_slide(prs, 1, "Problem Statement")

    _content_slide(
        prs,
        "The Challenge",
        [
            "Adobe Cloud Manager pipelines run complex AEM deployments with build, security test, and deploy stages — failures are frequent and costly.",
            "When a pipeline fails, engineers face 10,000+ line log files with no clear signal on root cause.",
            "Pre-deployment risk is invisible: teams promote code without knowing if similar changes failed before.",
            "Failure knowledge is siloed — past incidents aren't searchable or reusable across teams.",
            "Manual triage consumes hours per failure; recurring patterns go undetected for weeks.",
            "Business impact: delayed releases, wasted compute, frustrated developers, and SLA breaches.",
        ],
        subtitle="Why IDFC First Bank needed an intelligence layer on top of Cloud Manager",
    )

    _two_column_slide(
        prs,
        "Pain Points vs. Goals",
        "Current Pain Points",
        [
            "Log noise overwhelms human analysis",
            "No pre-deploy risk scoring from code changes",
            "Failures repeat because fixes aren't institutionalized",
            "Cross-deployment diffs require manual Splunk queries",
            "Code-to-log correlation is manual and slow",
            "No unified view of pipeline health",
        ],
        "Target Outcomes",
        [
            "Automated root cause analysis in minutes",
            "Evidence-weighted risk before every promotion",
            "Searchable memory of past failures & fixes",
            "Side-by-side execution comparison with AI narrative",
            "Pinpoint exact file/line that caused a failure",
            "Executive dashboard with KPIs and trends",
        ],
    )

    # ── Section 2: Solution ──
    _section_slide(prs, 2, "Solution Overview")

    _content_slide(
        prs,
        "What Is DevOps Agent 2?",
        [
            "A clean rebuild of the Adobe Cloud Manager pipeline intelligence agent — purpose-built for IDFC First Bank (Program 19905).",
            "Ingests execution metadata (Splunk), raw logs (Azure File Share), and Git diffs (Cloud Manager Git).",
            "Compresses noisy data through an 8-layer pipeline before sending focused context to LLMs.",
            "Returns structured Pydantic reports — not free-form text — for dashboard rendering and RAG indexing.",
            "Works offline with checked-in Splunk CSV exports; goes live with Splunk, Azure, Git, and LLM credentials.",
            "Dual interface: CLI for automation/CI and Streamlit dashboard for interactive analysis.",
        ],
    )

    _content_slide(
        prs,
        "Value Proposition",
        [
            ("⏱ Time to Resolution", 0),
            ("Reduce failure triage from hours to minutes via LogSage log pruning + AI RCA", 1),
            ("🛡 Proactive Risk Prevention", 0),
            ("Score every commit before promotion — rule engine + LLM + optional ML model", 1),
            ("🧠 Institutional Memory", 0),
            ("ChromaDB vector store + hybrid retrieval finds similar past incidents automatically", 1),
            ("📊 Executive Visibility", 0),
            ("30-day failure reports with estimated hours wasted and prioritized action items", 1),
            ("🔒 Enterprise Ready", 0),
            ("Secret redaction, credential-safe Git connector, structured JSON output, offline smoke tests", 1),
        ],
    )

    # ── Section 3: Architecture ──
    _section_slide(prs, 3, "Architecture")

    _content_slide(
        prs,
        "8-Layer Data Processing Pipeline",
        [
            "Layer 1 — Collection: Splunk metadata, Azure logs (build/security/deploy), Cloud Manager Git",
            "Layer 2 — Parsing: 10K log lines → structured error_type, error_message, key_lines",
            "Layer 3 — Enrichment: Link execution → commit SHA → changed files → probable cause",
            "Layer 4 — Compression: 99 executions → top 5 patterns; ~80% token reduction",
            "Layer 5 — Context Assembly: Feature-specific bundles (report, risk, pinpoint, scan)",
            "Layer 6 — Prompt Engineering: Role + JSON schema + constraints (no invented numbers)",
            "Layer 7 — LLM Call: Per-feature temperature/token tuning (Azure OpenAI / Anthropic / Gemini)",
            "Layer 8 — Output Parsing: Pydantic validation with auto-retry on schema failure",
        ],
        subtitle="Every analysis flows through the same disciplined pipeline",
    )

    _content_slide(
        prs,
        "Data Sources & Connectors",
        [
            ("Splunk Connector", 0),
            ("Live API or CSV fallback — pipeline list, failed steps, share name mapping", 1),
            ("Azure File Share Connector", 0),
            ("Fetches build.log, securityTests.log, deploy.log per execution", 1),
            ("Cloud Manager Git Connector", 0),
            ("Commit diffs, changed files, AEM module detection — no credential persistence in URLs", 1),
            ("Cloud Manager API Connector", 0),
            ("Resolves execution → commit SHA; multi-tenant token support (IDFC, HDFC, Apollo)", 1),
            ("ML Service Client", 0),
            ("Optional devops-risk-ml predictions with SHAP feature contributions", 1),
        ],
    )

    # ── Section 4: Features ──
    _section_slide(prs, 4, "Core Features")

    _table_slide(
        prs,
        "Six Intelligence Capabilities",
        ["Feature", "CLI Command", "Purpose"],
        [
            ["1. Failure Analysis Report", "report", "30-day executive summary with recurring patterns & fixes"],
            ["2. Pre-Deploy Risk Assessment", "risk --commit <sha>", "Score build/security/deploy risk from code changes"],
            ["3. Cross-Deployment Comparison", "compare --exec-a / --exec-b", "Diff two executions with AI narrative"],
            ["4. Code-to-Log Correlation", "correlate --execution-id", "Map parsed log errors to repo files"],
            ["5a. Proactive Code Scan", "scan", "Find antipatterns before they break pipelines"],
            ["5b. Reactive Failure Pinpoint", "pinpoint --execution-id", "Exact file + line that caused failure"],
            ["6. Post-Failure Assessment", "assess-failure --execution-id", "LogSage two-stage RCA + risk + retry advice"],
        ],
    )

    _content_slide(
        prs,
        "Feature 1: Failure Analysis Report",
        [
            "Analyzes last 30 days of pipeline executions from Splunk data.",
            "Identifies critical and recurring failure patterns by step and error type.",
            "Estimates hours wasted on failed/cancelled runs.",
            "Produces executive summary bullets, root cause narratives, and prioritized fix recommendations.",
            "Output: FailureReport Pydantic model → Markdown + JSON → indexed in ChromaDB for RAG.",
            "CLI: python3 cli.py report [--no-llm] [--no-logs]",
        ],
    )

    _content_slide(
        prs,
        "Feature 2: Pre-Deployment Risk Assessment",
        [
            "Evidence-weighted scoring based on WHAT CHANGED in the commit — not raw failure counts.",
            "CommitProfile analysis: AEM modules touched, dependency changes, config edits, anti-patterns.",
            "Per-step risk levels: build · securityTest · deploy with causal rationale.",
            "Technical failure hypotheses with likelihood, confidence, and verification steps.",
            "Blast radius analysis: affected modules, rollback complexity, user-facing impact.",
            "Optional ML overlay: overall risk score, step probabilities, GO/HOLD/NO_GO recommendation.",
            "CLI: python3 cli.py risk --commit <sha>",
        ],
    )

    _content_slide(
        prs,
        "Features 3–5: Compare, Correlate, Scan & Pinpoint",
        [
            ("Cross-Deployment Comparison", 0),
            ("Side-by-side two executions: step outcomes, log diffs, commit changes, AI summary", 1),
            ("Code-to-Log Correlation", 0),
            ("Given a failed execution (or parsed error JSON), resolve relevant code snippets via file resolver", 1),
            ("Proactive Code Scan", 0),
            ("Scan AEM repo for Cloud Manager antipatterns: OSGi configs, dispatcher rules, pom.xml issues", 1),
            ("Reactive Failure Pinpoint", 0),
            ("Parse failure log → grep repo → LLM ranks exact file:line culprits with confidence scores", 1),
        ],
    )

    # ── Section 5: LogSage ──
    _section_slide(prs, 5, "LogSage Pipeline")

    _content_slide(
        prs,
        "LogSage: Two-Stage Post-Failure Assessment",
        [
            "Stage 1 — Log Processing (no LLM):",
            ("Drain3 template extraction from success logs for noise baseline", 1),
            ("Filter anomalous lines vs. success templates", 1),
            ("Expand context windows around error candidates", 1),
            ("Token pruning to fit ~22K token budget (tiktoken)", 1),
            "Stage 1 RCA — LLM extracts error_type, root_cause, cascading failures from pruned logs",
            "Stage 2 — Hybrid retrieval finds similar incidents + LLM produces PostFailureRiskReport",
            "Outputs: retry recommendation (RETRY_SAFE / RETRY_WITH_FIX / DO_NOT_RETRY), fix steps, blast radius",
        ],
    )

    _content_slide(
        prs,
        "LogSage Components",
        [
            "drain_templates.py — Drain3 log template database from successful runs",
            "log_filter.py — Anomaly detection against success baseline",
            "context_window.py — Expand ±N lines around error candidates",
            "token_pruner.py — Greedy pruning to stay within LOGSAGE_TOKEN_LIMIT",
            "success_log_fetcher.py — Pull recent successful logs for same pipeline/step",
            "post_failure_assessor.py — Orchestrates Stage 1 → RCA → retrieval → Stage 2",
        ],
        subtitle="Designed to handle 10,000+ line logs within LLM context limits",
    )

    # ── Section 6: RAG ──
    _section_slide(prs, 6, "Retrieval & Memory")

    _content_slide(
        prs,
        "Hybrid Retrieval (8 Routes)",
        [
            "Query normalization with Jaccard deduplication and token capping",
            "Query rewriting — multiple search variants from RCA output",
            "HyDE — hypothetical document embedding for semantic search",
            "BM25 lexical search over indexed failure records",
            "Relational route — match step + error_type in failure history",
            "Vector search via ChromaDB embeddings",
            "Candidate merge, dedup, and BGE reranker (bge-reranker-v2-m3) for final ranking",
            "Every analysis output is stored back → becomes RAG memory for future queries",
        ],
    )

    _content_slide(
        prs,
        "Institutional Memory (ChromaDB)",
        [
            "Vector store indexes failure reports, scan findings, and post-failure assessments.",
            "Dashboard pages: Memory Search (semantic query) and Memory Explorer (browse clusters).",
            "3D PCA visualization of failure embeddings — clusters reveal semantic similarity.",
            "Memory stats shown on Overview dashboard: failure count, scan findings, cache TTL.",
            "Enables: 'Have we seen this error before?' with similarity scores and past fixes.",
        ],
    )

    # ── Section 7: ML ──
    _section_slide(prs, 7, "ML Integration")

    _content_slide(
        prs,
        "ML Risk Service Integration",
        [
            "Optional devops-risk-ml microservice at ML_SERVICE_URL (default localhost:8090).",
            "Predicts overall_risk_score (0–1) and per-step probabilities from commit features.",
            "SHAP top_features explain which commit signals drive the score.",
            "Promotion recommendation: GO · HOLD · NO_GO overlaid on LLM risk report.",
            "Resolves commit SHA from dev execution ID via Cloud Manager API.",
            "Graceful degradation — platform works fully without ML service.",
        ],
    )

    # ── Section 8: UI ──
    _section_slide(prs, 8, "User Interfaces")

    _content_slide(
        prs,
        "Streamlit Dashboard (8 Pages)",
        [
            "Overview — KPIs, failure-by-step chart, status pie, 30-day trend, recent failures",
            "Failure Analysis — Run/view AI failure report with critical & recurring findings",
            "Risk Assessment — Select commit, run pre-deploy risk with step breakdown & ML score",
            "Failure Pinpoint — Pick failed execution, pinpoint code culprits interactively",
            "Pipeline Intelligence — Deep execution explorer with log viewer",
            "Static Analysis — Proactive repo scan for antipatterns",
            "Memory Search — Semantic search over indexed failures",
            "Memory Explorer — Browse and visualize ChromaDB memory clusters",
        ],
        subtitle="streamlit run dashboard/app.py",
    )

    _content_slide(
        prs,
        "CLI Entry Points",
        [
            "python3 cli.py report [--no-llm] [--no-logs]",
            "python3 cli.py risk --commit <sha>",
            "python3 cli.py compare --exec-a <id> --exec-b <id>",
            "python3 cli.py correlate --execution-id <id>",
            "python3 cli.py scan",
            "python3 cli.py pinpoint --execution-id <id>",
            "python3 cli.py assess-failure --execution-id <id> [--no-reranker]",
            "Global flags: --no-llm (rules/data only), --no-logs (skip Azure fetch)",
        ],
    )

    # ── Section 9: Tech ──
    _section_slide(prs, 9, "Technology & Status")

    _two_column_slide(
        prs,
        "Technology Stack",
        "Core",
        [
            "Python 3 · Pydantic v2 models",
            "Azure OpenAI / Anthropic / Gemini",
            "ChromaDB vector store",
            "Streamlit dashboard + Plotly charts",
            "pandas · scikit-learn · tiktoken",
        ],
        "Specialized",
        [
            "Drain3 — log template mining",
            "rank-bm25 — lexical retrieval",
            "FlagEmbedding — BGE reranker",
            "azure-storage-file-share",
            "pytest offline smoke tests",
        ],
    )

    _content_slide(
        prs,
        "Current Status & Offline Demo",
        [
            "✅ Offline path works with checked-in Splunk CSV exports — no credentials required.",
            "✅ Smoke test: python3 -B tests/test_offline_smoke.py",
            "✅ Quick report: python3 -B cli.py report --no-llm --no-logs",
            "✅ Unit tests: LogSage filter, token pruner, hybrid retrieval",
            "🔧 Live mode requires: Splunk, Azure, Cloud Manager Git, LLM API keys in .env",
            "🔧 ML service optional for enhanced risk scoring",
            "📁 Reports saved to reports/ · Cache in data/cache/ · Vectors in data/chroma_db/",
        ],
    )

    _content_slide(
        prs,
        "Structured Output Models",
        [
            "FailureReport — executive summary, critical/recurring findings, hours wasted",
            "RiskReport — step risks, hypotheses, blast radius, ML prediction, evidence items",
            "PostFailureRiskReport — retry recommendation, fix steps, similar incidents",
            "LogSageRCAReport — Stage 1 intermediate RCA with error line refs",
            "PinpointReport / ScanReport — code-level findings with confidence",
            "All models: Pydantic-validated → JSON + Markdown → ChromaDB indexed",
        ],
    )

    # ── Closing ──
    _section_slide(prs, 10, "Summary")

    _content_slide(
        prs,
        "Key Takeaways",
        [
            "DevOps Agent 2 transforms Adobe Cloud Manager from a black box into an intelligence platform.",
            "8-layer pipeline ensures signal over noise — ~80% token reduction before every LLM call.",
            "Six capabilities cover the full lifecycle: report → risk → compare → correlate → scan → assess.",
            "LogSage solves the log volume problem with Drain3 + token pruning + two-stage RCA.",
            "Hybrid retrieval + ChromaDB memory means every failure makes the system smarter.",
            "Works today offline; scales to live Splunk/Azure/Git with credential configuration.",
        ],
    )

    _title_slide(
        prs,
        "Thank You",
        "DevOps Intelligence Platform · IDFC First Bank · Program 19905",
        "Questions & Demo: python3 cli.py report --no-llm --no-logs | streamlit run dashboard/app.py",
    )

    return prs


def main():
    out = Path(__file__).resolve().parent.parent / "docs" / "DevOps_Agent_2_Project_Overview.pptx"
    out.parent.mkdir(parents=True, exist_ok=True)
    prs = build_presentation()
    prs.save(str(out))
    print(f"Saved: {out}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
